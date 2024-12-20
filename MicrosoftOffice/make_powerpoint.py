#!/usr/bin/env python3
"""Generate a PowerPoint presentation from a CSV of data."""

import argparse
import csv
import getpass

from datetime import datetime
from enum import IntEnum
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Length


# Per https://python-pptx.readthedocs.io/en/latest/user/slides.html, this slide
# ordering layout is conventional but can be overridden by individual
# templates.
class Layout(IntEnum):
    """Simple enumeration of PowerPoint layouts."""

    TITLE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE_ONLY = 5
    BLANK = 6
    CONTENT_WITH_CAPTION = 7
    PICTURE_WITH_CAPTION = 8


def create_presentation_from_csv(
    csv_file: Path, template: Optional[Path] = None, outname: Optional[str] = None
):
    """Create a PowerPoint presentation from a CSV of data."""
    # Create a PowerPoint presentation object
    presentation = Presentation(str(template) if template else None)

    # Add a title slide
    title_slide = presentation.slides.add_slide(
        presentation.slide_layouts[Layout.TITLE]
    )

    title_slide.shapes.title.text = "Groundbreaking Results"
    title_slide.shapes.placeholders[
        1
    ].text = f"{getpass.getuser()}\t\t{datetime.now().strftime('%A, %B %d, %Y')}"

    # Read the CSV file
    with open(csv_file, "r", newline="", encoding="utf-8") as file:
        reader = csv.DictReader(file)
        rows = list(reader)

    # Insert a slide with a summary table
    table_slide = presentation.slides.add_slide(
        presentation.slide_layouts[Layout.BLANK]
    )

    # The table is 1/4 the size of the slide and centered
    table = table_slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(rows[0]),
        left=Length(presentation.slide_width / 4),
        top=Length(presentation.slide_height / 4),
        width=Length(presentation.slide_width / 2),
        height=Length(presentation.slide_height / 2),
    ).table

    # Add the table column headings
    for col_index, key in enumerate(rows[0].keys()):
        table.cell(0, col_index).text = key

    for row_index, row in enumerate(rows, start=1):
        # Add the row into the table slide
        for col_index, value in enumerate(row.values()):
            table.cell(row_index, col_index).text = value

        # Add a slide for each entry
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[Layout.TITLE_AND_CONTENT]
        )

        # Set slide title
        slide.shapes.title.text = f"{row['First Name']} {row['Last Name']}"

        # Set slide content
        score = row["Score"]
        slide.placeholders[1].text = f"Age: {row['Age']}\nScore: {score}"

        # This would be easier if there were a ChartPlaceholder already in
        # place, but I was too lazy to add it to the template
        chart_data = ChartData()
        chart_data.categories = ["Score", "Anti-Score"]
        chart_data.add_series(row["First Name"], (score, 10 - int(score)))

        # The chart is 1/4 the size of the slide, positioned with the
        # upper-left corner 1/3 of the lay along the upper-left/lower-right
        # diagonal.
        slide.shapes.add_chart(
            chart_type=XL_CHART_TYPE.PIE,
            x=Length(presentation.slide_width / 3),
            y=Length(presentation.slide_height / 3),
            cx=Length(presentation.slide_width / 2),
            cy=Length(presentation.slide_height / 2),
            chart_data=chart_data,
        )

    # Save the presentation
    output_path = csv_file.with_suffix(".pptx")
    if outname:
        output_path = output_path.with_stem(outname)

    presentation.save(str(output_path))
    print(f"Saved presentation as {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("csvfile", type=Path)
    parser.add_argument("--template", type=Path)
    parser.add_argument("--outname", type=str)

    args = parser.parse_args()
    create_presentation_from_csv(args.csvfile, args.template, args.outname)
